#!/usr/bin/env python
"""
generate_kit.py
Standalone script to generate a digital kit (PDF, PPTX, and app.json) for a street vendor.
This script reuses the existing RAG pipeline and QR code generator from the backend.

Usage:
    python backend/generate_kit.py --vendor-name "Ramu Chai Wala" --business-type "Tea Stall" --location "MG Road,Kumarpalya Layout, Bangalore" --upi-id "ramuchai@upi"

Dependencies:
    - reportlab
    - python-pptx
    (The other dependencies are already in the project: fastapi, ibm-watsonx-ai, chromadb, qrcode, Pillow, etc.)

Note: This script must be run from the project root or with the backend directory in the Python path.
"""

import os
import sys
import json
import argparse
from pathlib import Path

# Ensure the backend directory is in the path so we can import from it
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, BASE_DIR)
sys.path.insert(0, os.path.join(BASE_DIR, 'backend'))

from backend.rag_pipeline import answer, is_index_ready
from backend.qr_generator import generate_qr, generate_business_card

# For PDF generation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# For PPTX generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Constants
OUTPUT_DIR = os.path.join(BASE_DIR, "static", "generated")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# PDF generation function
def generate_pdf(vendor_name, business_type, location, upi_id, answer_text, qr_image_path, output_path):
    """
    Generate a PDF document containing the vendor's digital kit.
    """
    doc = SimpleDocTemplate(output_path, pagesize=LETTER,
                            rightMargin=30, leftMargin=30,
                            topMargin=30, bottomMargin=18)
    styles = getSampleStyleSheet()
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        alignment=TA_CENTER,
        textColor=colors.darkblue
    )
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        spaceAfter=12,
        textColor=colors.darkgreen
    )
    normal_style = styles['Normal']
    normal_style.fontSize = 12
    normal_style.spaceAfter = 6

    flowables = []

    # Title
    flowables.append(Paragraph(f"Digital Kit for {vendor_name}", title_style))
    flowables.append(Spacer(1, 12))

    # Vendor Info Table
    data = [
        ["Vendor Name:", vendor_name],
        ["Business Type:", business_type],
        ["Location:", location],
    ]
    if upi_id:
        data.append(["UPI ID:", upi_id])
    t = Table(data, colWidths=[2*inch, 4*inch])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.whitesmoke),
        ('TEXTCOLOR', (0,0), (-1,-1), colors.black),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('FONTNAME', (0,0), (-1,-1), 'Helvetica'),
        ('FONTSIZE', (0,0), (-1,-1), 10),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 1, colors.grey),
    ]))
    flowables.append(t)
    flowables.append(Spacer(1, 20))

    # Answer section
    flowables.append(Paragraph("Digital Advice Summary:", heading_style))
    # Split answer into paragraphs (by newline) and add each
    for para in answer_text.split('\n'):
        if para.strip():
            flowables.append(Paragraph(para, normal_style))
            flowables.append(Spacer(1, 6))
    flowables.append(Spacer(1, 20))

    # QR Code Image
    if qr_image_path and os.path.exists(qr_image_path):
        flowables.append(Paragraph("UPI QR Code:", heading_style))
        img = Image(qr_image_path, width=2*inch, height=2*inch)
        flowables.append(img)
        flowables.append(Spacer(1, 20))

    # Build PDF
    doc.build(flowables)
    return output_path

# PPTX generation function
def generate_pptx(vendor_name, business_type, location, upi_id, answer_text, qr_image_path, output_path):
    """
    Generate a PowerPoint presentation containing the vendor's digital kit.
    """
    prs = Presentation()
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Digital Kit for {vendor_name}"
    subtitle.text = f"{business_type}\n{location}"
    if upi_id:
        subtitle.text += f"\nUPI: {upi_id}"

    # Slide 2: QR Code
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_shape.text_frame
    tf.text = "UPI QR Code"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Image
    if qr_image_path and os.path.exists(qr_image_path):
        left = Inches(3)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(qr_image_path, left, top, height=Inches(3))
    else:
        # Placeholder text if no image
        left = Inches(3)
        top = Inches(1.5)
        shape = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
        tf = shape.text_frame
        tf.text = "QR Code Image Not Available"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 3: Advice (split into multiple slides if needed)
    # We'll break the answer into bullet points (by line) and put multiple bullets per slide
    lines = [line.strip() for line in answer_text.split('\n') if line.strip()]
    # Simple chunking: 5 bullets per slide
    chunk_size = 5
    for i in range(0, len(lines), chunk_size):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Digital Advice"
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()  # Clear existing paragraph
        for j, line in enumerate(lines[i:i+chunk_size]):
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.level = 0
            p.font.size = Pt(18)
        # If we want to keep formatting simple, we can just set the whole text
        # But the above loop adds paragraphs with bullets.

    # Slide 4: Footer / Contact
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Contact & Support"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = (
        f"Generated by Street Vendor Digitalization Agent\n"
        f"Powered by IBM watsonx.ai\n"
        f"For assistance, contact: support@example.com"
    )
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(output_path)
    return output_path

# app.json generation function
def generate_app_json(vendor_name, business_type, location, upi_id, qr_image_url, answer_text, output_path):
    """
    Generate a simple app.json configuration file for a hypothetical vendor app.
    """
    app_config = {
        "appName": f"{vendor_name} Digital Kit",
        "version": "1.0.0",
        "vendor": {
            "name": vendor_name,
            "businessType": business_type,
            "location": location,
            "upiId": upi_id if upi_id else None
        },
        "resources": {
            "qrCode": qr_image_url  # This should be a URL accessible by the app
            # We could also include a link to the PDF/PPTX if hosted
        },
        "content": {
            "advice": answer_text
            # We could break the answer into sections, but for simplicity we put the whole text
        },
        "theme": {
            "primaryColor": "#1B2A6B",   # INDIGO from qr_generator
            "secondaryColor": "#F5A623", # TURMERIC
            "backgroundColor": "#FDF6E3" # KRAFT_LIGHT
        }
    }
    with open(output_path, 'w') as f:
        json.dump(app_config, f, indent=2)
    return output_path

def main():
    parser = argparse.ArgumentParser(description="Generate digital kit (PDF, PPTX, app.json) for a street vendor.")
    parser.add_argument("--vendor-name", required=True, help="Name of the vendor")
    parser.add_argument("--business-type", required=True, help="Type of business (e.g., Tea Stall, Vegetable Vendor)")
    parser.add_argument("--location", required=True, help="Location of the vendor")
    parser.add_argument("--upi-id", help="UPI ID for payments (optional)")
    parser.add_argument("--language", default="en", help="Language for the advice (default: en)")
    parser.add_argument("--output-dir", default=OUTPUT_DIR, help=f"Directory to save output files (default: {OUTPUT_DIR})")
    args = parser.parse_args()

    # Ensure the index is ready (we'll try to use the RAG pipeline)
    if not is_index_ready():
        print("Warning: Vector index not ready. Attempting to build index...")
        # We could call build_index here, but that might be heavy. Instead, we'll rely on the existing index.
        # For simplicity, we'll just proceed and let the answer function handle errors.
        pass

    # Generate the advice using the RAG pipeline
    # We'll create a query similar to the one in /api/generate-kit
    enriched_query = (
        f"I am {args.vendor_name}, a {args.business_type} located at {args.location}. "
        f"My UPI ID is {args.upi_id or 'not set yet'}. "
        f"How do I go digital? What schemes am I eligible for? "
        f"Which platforms should I list on? What are my local SEO tips?"
    )
    try:
        result = answer(enriched_query, top_k=3, language=args.language)
        answer_text = result["answer"]
        # We don't use the retrieved docs in the output files for simplicity, but we could.
    except Exception as e:
        print(f"Error generating advice: {e}")
        # Fallback to a generic message
        answer_text = (
            f"Digital advice for {args.vendor_name}:\n"
            "- Embrace digital payments (UPI, QR codes).\n"
            "- List your business on Google My Business and Justdial.\n"
            "- Use social media for promotion.\n"
            "- Explore government schemes for street vendors.\n"
            "- Keep digital records of sales and inventory."
        )

    # Generate QR code image (we'll generate the business card image which includes QR)
    # Alternatively, we could generate just the QR code. Let's generate the business card for a richer image.
    try:
        # Use the existing function to generate a business card PNG
        card_url = generate_business_card(
            vendor_name=args.vendor_name,
            business_type=args.business_type,
            location=args.location,
            upi_id=args.upi_id if args.upi_id else "dummy@upi"  # The function expects a UPI ID; if none, we use a dummy.
        )
        # The function returns a URL path like "/static/generated/card_xxxx.png"
        # We need the absolute file path
        qr_image_path = os.path.join(BASE_DIR, "static", "generated", os.path.basename(card_url))
        # Also keep the URL for app.json
        qr_image_url = card_url  # This is relative to the server root
    except Exception as e:
        print(f"Error generating business card: {e}")
        # Fallback to just QR code
        try:
            qr_url = generate_qr(args.upi_id or "dummy@upi", args.vendor_name)
            qr_image_path = os.path.join(BASE_DIR, "static", "generated", os.path.basename(qr_url))
            qr_image_url = qr_url
        except Exception as e2:
            print(f"Error generating QR code: {e2}")
            qr_image_path = None
            qr_image_url = None

    # Define output file paths
    pdf_filename = f"kit_{args.vendor_name.replace(' ', '_')}.pdf"
    pptx_filename = f"kit_{args.vendor_name.replace(' ', '_')}.pptx"
    appjson_filename = f"app_{args.vendor_name.replace(' ', '_')}.json"

    pdf_path = os.path.join(args.output_dir, pdf_filename)
    pptx_path = os.path.join(args.output_dir, pptx_filename)
    appjson_path = os.path.join(args.output_dir, appjson_filename)

    # Generate PDF
    if qr_image_path:
        print(f"Generating PDF: {pdf_path}")
        generate_pdf(
            vendor_name=args.vendor_name,
            business_type=args.business_type,
            location=args.location,
            upi_id=args.upi_id,
            answer_text=answer_text,
            qr_image_path=qr_image_path,
            output_path=pdf_path
        )
    else:
        print("Skipping PDF generation because QR image is not available.")

    # Generate PPTX
    if qr_image_path:
        print(f"Generating PPTX: {pptx_path}")
        generate_pptx(
            vendor_name=args.vendor_name,
            business_type=args.business_type,
            location=args.location,
            upi_id=args.upi_id,
            answer_text=answer_text,
            qr_image_path=qr_image_path,
            output_path=pptx_path
        )
    else:
        print("Skipping PPTX generation because QR image is not available.")

    # Generate app.json
    print(f"Generating app.json: {appjson_path}")
    generate_app_json(
        vendor_name=args.vendor_name,
        business_type=args.business_type,
        location=args.location,
        upi_id=args.upi_id,
        qr_image_url=qr_image_url,
        answer_text=answer_text,
        output_path=appjson_path
    )

    print("Generation complete.")

if __name__ == "__main__":
    main()